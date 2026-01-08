import os
import requests
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()


def get_stock_data(type_key):
    """Generic fetcher for gainers or losers."""
    url = os.getenv(f'TOP_{type_key.upper()}_URL')
    api_key = os.getenv('FMP_API_KEY')

    if not url or not api_key:
        raise ValueError(f"Environment variables for {type_key} or API Key missing")

    try:
        response = requests.get(f"{url}?apikey={api_key}", timeout=15)
        response.raise_for_status()
        data = response.json()

        # Save raw data
        with open(f'top_{type_key}_raw.json', 'w') as f:
            json.dump(data, f, indent=2)

        return data
    except Exception as e:
        print(f"Error fetching {type_key}: {e}")
        return []


def process_data(data, type_key):
    """Filters and formats data efficiently using Pandas."""
    if not data:
        return pd.DataFrame()

    df = pd.DataFrame(data)
    filter_pattern = r'^[a-zA-Z]{1,4}$'

    mask = (
            (df['exchange'].isin(['NASDAQ', 'NYSE'])) &
            (df['symbol'].str.match(filter_pattern)) &
            (df['price'] >= 50.00)
    )

    if type_key == 'gainers':
        mask &= (df['changesPercentage'] >= 10.00)
        sort_asc = False
    else:
        mask &= (df['changesPercentage'] <= -10.00)
        sort_asc = True

    df = df[mask].sort_values('changesPercentage', ascending=sort_asc).head(10).copy()

    # Vectorized formatting
    df['changesPercentage'] = df['changesPercentage'].map("{:.1f}%".format)
    df['price'] = df['price'].map("${:.2f}".format)

    return df[['symbol', 'name', 'price', 'changesPercentage', 'exchange']]


def create_split_slide(gainers_df, losers_df):
    """Creates a single 16:9 slide with side-by-side tables in brand colors."""
    prs = Presentation()
    # Setting dimensions for 16:9 at 144 DPI (Inches * 144 = Pixels)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank Slide

    # Define Brand Colors
    GAINER_COLOR = RGBColor(67, 4, 183)  # #4304B7
    LOSER_COLOR = RGBColor(99, 46, 98)  # #632E62

    def add_table_to_side(df, title_text, left_pos, brand_color):
        # Section Title
        txBox = slide.shapes.add_textbox(left_pos, Inches(0.4), Inches(6.2), Inches(0.6))
        tf = txBox.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = brand_color
        p.alignment = PP_ALIGN.CENTER

        # Table Dimensions
        rows, cols = len(df) + 1, len(df.columns)
        table_shape = slide.shapes.add_table(rows, cols, left_pos, Inches(1.2), Inches(6.2), Inches(5.0))
        table = table_shape.table

        # Column Widths (Adjusted for the split view)
        table.columns[0].width = Inches(0.8)  # Symbol
        table.columns[1].width = Inches(2.4)  # Name
        table.columns[2].width = Inches(1.0)  # Price
        table.columns[3].width = Inches(1.0)  # Change
        table.columns[4].width = Inches(1.0)  # Exchange

        # Format Headers
        headers = ['Ticker', 'Company Name', 'Price', 'Chg%', 'Exchange']
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = brand_color

            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.size = Pt(12)
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER

        # Format Rows
        for r_idx, (_, row) in enumerate(df.iterrows(), 1):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)

                # Zebra striping
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

                # Alignment
                if c_idx in [0, 4]:
                    para.alignment = PP_ALIGN.CENTER
                elif c_idx in [2, 3]:
                    para.alignment = PP_ALIGN.RIGHT
                else:
                    para.alignment = PP_ALIGN.LEFT

    # Add Left Table (Gainers)
    add_table_to_side(gainers_df, "Top Gainers", Inches(0.3), GAINER_COLOR)

    # Add Right Table (Losers)
    add_table_to_side(losers_df, "Top Losers", Inches(6.8), LOSER_COLOR)

    # Add Source at the very bottom
    source_box = slide.shapes.add_textbox(Inches(0), Inches(6.9), Inches(13.333), Inches(0.4))
    tf_source = source_box.text_frame
    tf_source.text = "Source: Prepared by Ki-Wealth based on FMP data"
    p_source = tf_source.paragraphs[0]
    p_source.alignment = PP_ALIGN.CENTER
    p_source.font.size = Pt(12)
    p_source.font.italic = True

    prs.save('top_gainers_and_losers_for_the_day.pptx')


def main():
    print("--- Starting Professional Report Generation ---")
    g_raw = get_stock_data('gainers')
    l_raw = get_stock_data('losers')

    g_df = process_data(g_raw, 'gainers')
    l_df = process_data(l_raw, 'losers')

    if g_df.empty and l_df.empty:
        print("No data matched the criteria today.")
        return

    create_split_slide(g_df, l_df)
    print("Success! One-slide presentation created with brand colors.")


if __name__ == "__main__":
    main()